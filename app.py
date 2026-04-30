from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io, os, json, requests, random

app = Flask(__name__)
CORS(app)

W = 9144000
H = 6858000

# ── Color palettes — topic ke hisaab se ──────────────────────
PALETTES = [
    {'primary': '7c3aed', 'secondary': 'a78bfa', 'bg': 'faf5ff', 'dark': '1a0533', 'body': '374151', 'light': 'ede9fe'},
    {'primary': '0f766e', 'secondary': '2dd4bf', 'bg': 'f0fdfa', 'dark': '042f2e', 'body': '374151', 'light': 'ccfbf1'},
    {'primary': 'b91c1c', 'secondary': 'f87171', 'bg': 'fff1f2', 'dark': '450a0a', 'body': '374151', 'light': 'fecdd3'},
    {'primary': '1d4ed8', 'secondary': '60a5fa', 'bg': 'eff6ff', 'dark': '1e3a8a', 'body': '374151', 'light': 'dbeafe'},
    {'primary': '15803d', 'secondary': '4ade80', 'bg': 'f0fdf4', 'dark': '052e16', 'body': '374151', 'light': 'bbf7d0'},
    {'primary': 'c2410c', 'secondary': 'fb923c', 'bg': 'fff7ed', 'dark': '431407', 'body': '374151', 'light': 'fed7aa'},
    {'primary': '6d28d9', 'secondary': 'c4b5fd', 'bg': 'f5f3ff', 'dark': '2e1065', 'body': '374151', 'light': 'ddd6fe'},
    {'primary': '0369a1', 'secondary': '38bdf8', 'bg': 'f0f9ff', 'dark': '082f49', 'body': '374151', 'light': 'bae6fd'},
    {'primary': '9f1239', 'secondary': 'fb7185', 'bg': 'fff1f2', 'dark': '4c0519', 'body': '374151', 'light': 'fecdd3'},
    {'primary': '065f46', 'secondary': '34d399', 'bg': 'ecfdf5', 'dark': '022c22', 'body': '374151', 'light': 'a7f3d0'},
]

THEME_MAP = {
    'purple': PALETTES[0], 'blue': PALETTES[3], 'green': PALETTES[4],
    'orange': PALETTES[5], 'dark': PALETTES[6],
}

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = hex_rgb(fill); s.line.fill.background()
    return s

def circle(slide, x, y, size, fill):
    s = slide.shapes.add_shape(9, Emu(x), Emu(y), Emu(size), Emu(size))
    s.fill.solid(); s.fill.fore_color.rgb = hex_rgb(fill); s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=18, bold=False, color='000000',
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = hex_rgb(color); r.font.name = 'Calibri'
    return tb

def bullets(slide, items, x, y, w, h, size=15, color='374151', prefix='• '):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run(); r.text = prefix + str(item)
        r.font.size = Pt(size); r.font.color.rgb = hex_rgb(color); r.font.name = 'Calibri'

# ═══════════════════════════════════════════════════════════════
# TITLE SLIDE — 4 different designs
# ═══════════════════════════════════════════════════════════════
def title_v1(slide, data, P):
    """Full dark background, big centered text"""
    rect(slide, 0, 0, W, H, P['dark'])
    circle(slide, W-2500000, -600000, 3200000, P['primary'])
    circle(slide, -800000, H-1500000, 2500000, P['secondary'] + '40')
    rect(slide, 0, H-200000, W, 200000, P['primary'])
    if data.get('emoji'): txt(slide, data['emoji'], W//2-500000, 1000000, 1000000, 1000000, size=56, align=PP_ALIGN.CENTER)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 600000, 2200000, W-1200000, 1500000, size=54, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
    if data.get('subtitle') or data.get('subheading'): txt(slide, data['subtitle'], 600000, 3900000, W-1200000, 700000, size=22, color=P['secondary'], italic=True, align=PP_ALIGN.CENTER)
    txt(slide, 'HackMate AI', 600000, H-500000, W-1200000, 350000, size=11, color='6b7280', align=PP_ALIGN.CENTER)

def title_v2(slide, data, P):
    """Split — left colored panel, right white"""
    rect(slide, 0, 0, W//2, H, P['primary'])
    rect(slide, W//2, 0, W//2, H, 'FFFFFF')
    circle(slide, W//2-1200000, H//2-1200000, 2400000, P['secondary'] + '30')
    if data.get('emoji'): txt(slide, data['emoji'], 200000, 600000, 900000, 900000, size=52)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 200000, 1700000, W//2-400000, 1800000, size=42, bold=True, color='FFFFFF')
    if data.get('subtitle') or data.get('subheading'): txt(slide, data['subtitle'], 200000, 3700000, W//2-400000, 800000, size=18, color='E0D7FF', italic=True)
    txt(slide, 'HackMate AI', 200000, H-500000, W//2-400000, 350000, size=11, color='C4B5FD')
    # Right side decoration
    circle(slide, W//2+400000, 800000, 1200000, P['light'])
    circle(slide, W//2+1800000, 2000000, 800000, P['secondary']+'50')
    circle(slide, W//2+600000, 3200000, 1600000, P['light'])

def title_v3(slide, data, P):
    """Bottom strip design"""
    rect(slide, 0, 0, W, H, 'FFFFFF')
    rect(slide, 0, H-1800000, W, 1800000, P['primary'])
    rect(slide, 0, H-1820000, W, 40000, P['secondary'])
    circle(slide, 200000, 200000, 1200000, P['light'])
    circle(slide, W-1600000, 400000, 2000000, P['light'])
    if data.get('emoji'): txt(slide, data['emoji'], W//2-400000, 700000, 800000, 800000, size=48, align=PP_ALIGN.CENTER)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 400000, 1700000, W-800000, 1500000, size=52, bold=True, color=P['dark'], align=PP_ALIGN.CENTER)
    if data.get('subtitle') or data.get('subheading'): txt(slide, data['subtitle'], 400000, H-1600000, W-800000, 600000, size=20, color='FFFFFF', italic=True, align=PP_ALIGN.CENTER)
    txt(slide, 'HackMate AI', 400000, H-800000, W-800000, 350000, size=11, color='C4B5FD', align=PP_ALIGN.CENTER)

def title_v4(slide, data, P):
    """Diagonal design"""
    rect(slide, 0, 0, W, H, P['bg'])
    # Diagonal colored block (approximate with large rotated rect)
    rect(slide, 0, 0, W, H//2+500000, P['primary'])
    rect(slide, 0, H//2-200000, W, H//2+200000, P['dark'])
    # Circles
    circle(slide, W-1800000, 200000, 1400000, P['secondary']+'40')
    circle(slide, 400000, H//2+200000, 1000000, P['secondary']+'30')
    if data.get('emoji'): txt(slide, data['emoji'], 500000, 500000, 900000, 900000, size=50)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 500000, 1500000, W-1000000, 1500000, size=48, bold=True, color='FFFFFF')
    if data.get('subtitle') or data.get('subheading'): txt(slide, data['subtitle'], 500000, H//2+200000, W-1000000, 700000, size=20, color=P['secondary'], italic=True)
    txt(slide, 'HackMate AI', 500000, H-500000, W-1000000, 350000, size=11, color='6b7280')

TITLE_VARIANTS = [title_v1, title_v2, title_v3, title_v4]

# ═══════════════════════════════════════════════════════════════
# PROBLEM SLIDE — 4 variants
# ═══════════════════════════════════════════════════════════════
def problem_v1(slide, data, P):
    """Dark left panel + numbered cards"""
    rect(slide, 0, 0, 3000000, H, P['dark'])
    rect(slide, 3000000, 0, W-3000000, H, 'FFFFFF')
    if data.get('emoji'): txt(slide, data['emoji'], 200000, 400000, 800000, 800000, size=40)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 200000, 1300000, 2600000, 800000, size=28, bold=True, color='FFFFFF')
    if data.get('headline'): txt(slide, data['headline'], 200000, 2300000, 2600000, 1400000, size=15, color=P['secondary'], italic=True)
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    ch, gap = 1500000, 150000
    sy = (H-(len(items[:3])*ch+(len(items[:3])-1)*gap))//2
    for i, b in enumerate(items[:3]):
        y = sy + i*(ch+gap)
        rect(slide, 3200000, y, W-3500000, ch, P['light'])
        circle(slide, 3380000, y+400000, 650000, P['primary'])
        txt(slide, str(i+1), 3380000, y+400000, 650000, 650000, size=18, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
        txt(slide, b, 4180000, y+150000, W-4650000, ch-300000, size=14, color='374151')

def problem_v2(slide, data, P):
    """Top header + 3 horizontal stat bars"""
    rect(slide, 0, 0, W, 1400000, P['primary'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 280000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1100000, 320000, W-1400000, 700000, size=30, bold=True, color='FFFFFF')
    if data.get('headline'): txt(slide, data['headline'], 1100000, 980000, W-1400000, 370000, size=15, color='E0D7FF', italic=True)
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    icons = ['⚠️','❌','😰']
    for i, b in enumerate(items[:3]):
        y = 1700000 + i*1500000
        rect(slide, 400000, y, W-800000, 1200000, P['light'] if i%2==0 else 'FFFFFF')
        txt(slide, icons[i%3], 500000, y+300000, 600000, 600000, size=26)
        txt(slide, b, 1300000, y+200000, W-1900000, 800000, size=16, color='374151')

def problem_v3(slide, data, P):
    """Full dark bg + glowing cards"""
    rect(slide, 0, 0, W, H, P['dark'])
    circle(slide, -500000, -500000, 2000000, P['primary']+'30')
    circle(slide, W-1500000, H-1500000, 2000000, P['secondary']+'20')
    if data.get('emoji'): txt(slide, data['emoji'], 400000, 300000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1200000, 320000, W-1600000, 700000, size=30, bold=True, color='FFFFFF')
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    cw, gap = 2700000, 150000
    sx = (W-(3*cw+2*gap))//2
    for i, b in enumerate(items[:3]):
        x = sx + i*(cw+gap)
        rect(slide, x, 1400000, cw, H-1900000, P['primary']+'40')
        rect(slide, x, 1400000, cw, 8000, P['secondary'])
        txt(slide, ['⚠️','🔴','❗'][i], x+cw//2-300000, 1600000, 600000, 600000, size=28, align=PP_ALIGN.CENTER)
        txt(slide, b, x+150000, 2400000, cw-300000, H-2900000, size=14, color='FFFFFF', align=PP_ALIGN.CENTER)

def problem_v4(slide, data, P):
    """Timeline style"""
    rect(slide, 0, 0, W, H, 'FFFFFF')
    rect(slide, 0, 0, 120000, H, P['primary'])
    rect(slide, 0, 0, W, 1200000, P['light'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 250000, 650000, 650000, size=30)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1050000, 280000, W-1400000, 650000, size=28, bold=True, color=P['dark'])
    if data.get('headline'):
        txt(slide, f'"{data["headline"]}"', 300000, 1350000, W-600000, 500000, size=18, bold=True, color=P['primary'], italic=True)
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    for i, b in enumerate(items[:3]):
        y = 2100000 + i*1400000
        circle(slide, 250000, y+100000, 600000, P['primary'])
        txt(slide, str(i+1), 250000, y+100000, 600000, 600000, size=16, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
        rect(slide, 550000, y+350000, 60000, 1050000, P['light'])
        rect(slide, 1050000, y, W-1400000, 1100000, P['light'])
        txt(slide, b, 1200000, y+150000, W-1700000, 800000, size=15, color='374151')

PROBLEM_VARIANTS = [problem_v1, problem_v2, problem_v3, problem_v4]

# ═══════════════════════════════════════════════════════════════
# SOLUTION SLIDE — 4 variants
# ═══════════════════════════════════════════════════════════════
def solution_v1(slide, data, P):
    """Header + 3 cards in row"""
    rect(slide, 0, 0, W, 1400000, P['primary'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 250000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1100000, 300000, W-1500000, 700000, size=30, bold=True, color='FFFFFF')
    if data.get('headline'): txt(slide, data['headline'], 1100000, 970000, W-1500000, 370000, size=15, color='E0D7FF', italic=True)
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    icons = ['💡','⚡','🎯','🔥','✨']
    cw, gap = 2700000, 150000
    sx = (W-(3*cw+2*gap))//2
    for i, b in enumerate(items[:3]):
        x = sx + i*(cw+gap)
        rect(slide, x, 1600000, cw, H-2000000, P['light'])
        circle(slide, x+cw//2-350000, 1800000, 700000, P['primary'])
        txt(slide, icons[i%5], x+cw//2-350000, 1800000, 700000, 700000, size=24, align=PP_ALIGN.CENTER)
        txt(slide, b, x+150000, 2700000, cw-300000, H-3100000, size=14, color='374151', align=PP_ALIGN.CENTER)

def solution_v2(slide, data, P):
    """Dark bg + glowing feature pills"""
    rect(slide, 0, 0, W, H, P['dark'])
    circle(slide, W-1000000, -500000, 2000000, P['primary']+'40')
    if data.get('emoji'): txt(slide, data['emoji'], 400000, 300000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1200000, 320000, W-1600000, 700000, size=30, bold=True, color='FFFFFF')
    if data.get('headline'): txt(slide, data['headline'], 400000, 1200000, W-800000, 450000, size=18, bold=True, color=P['secondary'])
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    for i, b in enumerate(items[:4]):
        y = 1900000 + i*1100000
        rect(slide, 400000, y, W-800000, 900000, P['primary']+'35')
        rect(slide, 400000, y, 12000, 900000, P['secondary'])
        txt(slide, ['✅','🚀','⚡','🎯'][i%4], 600000, y+200000, 600000, 500000, size=22)
        txt(slide, b, 1350000, y+150000, W-1900000, 600000, size=16, color='FFFFFF')

def solution_v3(slide, data, P):
    """2+1 asymmetric grid"""
    rect(slide, 0, 0, W, H, 'FFFFFF')
    rect(slide, 0, 0, W, 1300000, P['light'])
    rect(slide, 0, 0, 120000, H, P['primary'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 250000, 650000, 650000, size=30)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1050000, 280000, W-1400000, 650000, size=28, bold=True, color=P['dark'])
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    # Card 1 big left
    rect(slide, 300000, 1500000, 4200000, H-2000000, P['primary'])
    txt(slide, '💡', 300000+400000, 1700000, 700000, 700000, size=32)
    if items: txt(slide, items[0], 300000+200000, 2600000, 3800000, H-3200000, size=16, color='FFFFFF')
    # Cards 2 & 3 stacked right
    for i in range(1, min(3, len(items))):
        y = 1500000 + (i-1)*((H-2000000)//2+100000)
        rect(slide, 4700000, y, W-5000000, (H-2200000)//2, P['light'])
        txt(slide, ['⚡','🎯'][i-1], 4850000, y+150000, 600000, 600000, size=26)
        txt(slide, items[i], 5600000, y+100000, W-5900000, (H-2200000)//2-200000, size=14, color='374151')

def solution_v4(slide, data, P):
    """Numbered list with big numbers"""
    rect(slide, 0, 0, W, H, P['bg'])
    rect(slide, 0, 0, W, 1300000, P['primary'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 250000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1100000, 300000, W-1500000, 700000, size=30, bold=True, color='FFFFFF')
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    for i, b in enumerate(items[:4]):
        y = 1500000 + i*1300000
        txt(slide, f'0{i+1}', 300000, y, 900000, 1100000, size=60, bold=True, color=P['primary']+'40')
        rect(slide, 1200000, y+200000, W-1600000, 900000, 'FFFFFF')
        txt(slide, b, 1400000, y+300000, W-2000000, 700000, size=16, color='374151')

SOLUTION_VARIANTS = [solution_v1, solution_v2, solution_v3, solution_v4]

# ═══════════════════════════════════════════════════════════════
# HOW IT WORKS — 4 variants
# ═══════════════════════════════════════════════════════════════
def how_v1(slide, data, P):
    """Horizontal steps with arrows"""
    rect(slide, 0, 0, W, 1200000, P['light'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 250000, 650000, 650000, size=30)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1050000, 280000, W-1400000, 650000, size=28, bold=True, color=P['dark'])
    steps = data.get('steps') or data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    n = min(len(steps), 4)
    if not n: return
    sw, aw = 1900000, 250000
    tw = n*sw+(n-1)*aw
    sx = (W-tw)//2
    for i, s in enumerate(steps[:4]):
        x = sx + i*(sw+aw)
        rect(slide, x, 1400000, sw, H-1900000, P['primary'])
        circle(slide, x+sw//2-350000, 1600000, 700000, 'FFFFFF')
        txt(slide, str(i+1), x+sw//2-350000, 1600000, 700000, 700000, size=22, bold=True, color=P['primary'], align=PP_ALIGN.CENTER)
        txt(slide, s, x+150000, 2500000, sw-300000, H-3000000, size=13, color='FFFFFF', align=PP_ALIGN.CENTER)
        if i < n-1: txt(slide, '→', x+sw+30000, 2900000, 190000, 400000, size=22, bold=True, color=P['primary'], align=PP_ALIGN.CENTER)

def how_v2(slide, data, P):
    """Vertical timeline"""
    rect(slide, 0, 0, W, H, 'FFFFFF')
    rect(slide, 0, 0, 120000, H, P['primary'])
    rect(slide, 0, 0, W, 1200000, P['light'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 250000, 650000, 650000, size=30)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1050000, 280000, W-1400000, 650000, size=28, bold=True, color=P['dark'])
    # Center vertical line
    rect(slide, W//2-30000, 1400000, 60000, H-1700000, P['light'])
    steps = data.get('steps') or data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    for i, s in enumerate(steps[:4]):
        y = 1500000 + i*1250000
        left = i % 2 == 0
        # Circle on center
        circle(slide, W//2-300000, y+100000, 600000, P['primary'])
        txt(slide, str(i+1), W//2-300000, y+100000, 600000, 600000, size=16, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
        if left:
            rect(slide, 300000, y, W//2-700000, 1050000, P['light'])
            txt(slide, s, 450000, y+150000, W//2-1000000, 750000, size=14, color='374151')
        else:
            rect(slide, W//2+400000, y, W//2-700000, 1050000, P['light'])
            txt(slide, s, W//2+550000, y+150000, W//2-1000000, 750000, size=14, color='374151')

def how_v3(slide, data, P):
    """Dark funnel style"""
    rect(slide, 0, 0, W, H, P['dark'])
    circle(slide, -400000, -400000, 1800000, P['primary']+'30')
    if data.get('emoji'): txt(slide, data['emoji'], 400000, 300000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1200000, 320000, W-1600000, 700000, size=30, bold=True, color='FFFFFF')
    steps = data.get('steps') or data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    for i, s in enumerate(steps[:4]):
        y = 1400000 + i*1300000
        w_box = W - 800000 - i*600000
        x_box = 400000 + i*300000
        rect(slide, x_box, y, w_box, 1000000, P['primary']+'50')
        rect(slide, x_box, y, 12000, 1000000, P['secondary'])
        circle(slide, x_box+80000, y+175000, 650000, P['secondary'])
        txt(slide, str(i+1), x_box+80000, y+175000, 650000, 650000, size=18, bold=True, color=P['dark'], align=PP_ALIGN.CENTER)
        txt(slide, s, x_box+900000, y+200000, w_box-1100000, 600000, size=15, color='FFFFFF')

def how_v4(slide, data, P):
    """2x2 grid steps"""
    rect(slide, 0, 0, W, H, P['bg'])
    rect(slide, 0, 0, W, 1300000, P['primary'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 280000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1100000, 320000, W-1500000, 650000, size=28, bold=True, color='FFFFFF')
    steps = data.get('steps') or data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    cw = (W-700000)//2
    ch = (H-1800000)//2
    gap = 200000
    for i, s in enumerate(steps[:4]):
        col, row = i%2, i//2
        x = 300000 + col*(cw+gap)
        y = 1500000 + row*(ch+gap)
        rect(slide, x, y, cw, ch, 'FFFFFF')
        rect(slide, x, y, cw, 10000, P['primary'])
        txt(slide, f'0{i+1}', x+100000, y+50000, 700000, 700000, size=36, bold=True, color=P['primary']+'50')
        txt(slide, s, x+150000, y+550000, cw-300000, ch-700000, size=14, color='374151')

HOW_VARIANTS = [how_v1, how_v2, how_v3, how_v4]

# ═══════════════════════════════════════════════════════════════
# TECH SLIDE — 3 variants
# ═══════════════════════════════════════════════════════════════
def tech_v1(slide, data, P):
    """Left bar + 2-col list"""
    rect(slide, 0, 0, 120000, H, P['primary'])
    rect(slide, 0, 0, W, 1200000, P['light'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 250000, 650000, 650000, size=30)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1050000, 280000, W-1400000, 650000, size=28, bold=True, color=P['dark'])
    rect(slide, 300000, 1200000, W-600000, 40000, P['primary'])
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    for i, item in enumerate(items[:8]):
        x = 400000 if i<4 else W//2+100000
        y = 1450000 + (i%4)*680000
        w = W//2-600000
        circle(slide, x, y+200000, 280000, P['primary'])
        txt(slide, item, x+400000, y, w, 680000, size=16, color='374151')

def tech_v2(slide, data, P):
    """Card grid"""
    rect(slide, 0, 0, W, H, P['dark'])
    if data.get('emoji'): txt(slide, data['emoji'], 400000, 300000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1200000, 320000, W-1600000, 700000, size=30, bold=True, color='FFFFFF')
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    cw = (W-700000)//2
    ch = 1100000
    gap = 150000
    for i, item in enumerate(items[:6]):
        col, row = i%2, i//2
        x = 300000 + col*(cw+gap)
        y = 1400000 + row*(ch+gap)
        rect(slide, x, y, cw, ch, P['primary']+'40')
        rect(slide, x, y, 12000, ch, P['secondary'])
        txt(slide, item, x+200000, y+300000, cw-400000, ch-400000, size=16, color='FFFFFF')

def tech_v3(slide, data, P):
    """Pill badges"""
    rect(slide, 0, 0, W, H, 'FFFFFF')
    rect(slide, 0, 0, W, 1300000, P['primary'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 280000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1100000, 320000, W-1500000, 650000, size=28, bold=True, color='FFFFFF')
    items = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    colors = [P['primary'], P['secondary'], '374151', '6d28d9', '0369a1', '065f46', '9f1239', '92400e']
    x, y = 300000, 1500000
    for i, item in enumerate(items[:8]):
        pill_w = max(len(item)*110000, 1800000)
        if x + pill_w > W - 300000:
            x = 300000; y += 1000000
        rect(slide, x, y, pill_w, 750000, colors[i%len(colors)])
        txt(slide, item, x+150000, y+150000, pill_w-300000, 450000, size=15, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
        x += pill_w + 200000

TECH_VARIANTS = [tech_v1, tech_v2, tech_v3]

# ═══════════════════════════════════════════════════════════════
# IMPACT SLIDE — 3 variants
# ═══════════════════════════════════════════════════════════════
def impact_v1(slide, data, P):
    """3 big stat cards"""
    rect(slide, 0, 0, W, H, P['light'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 280000, 650000, 650000, size=30)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1050000, 320000, W-1400000, 600000, size=28, bold=True, color=P['dark'])
    stats = data.get('stats', [])
    cw, ch, gap = 2700000, 3400000, 150000
    sx = (W-(3*cw+2*gap))//2
    for i, s in enumerate(stats[:3]):
        x = sx + i*(cw+gap)
        rect(slide, x, 1500000, cw, ch, P['primary'])
        txt(slide, s.get('number',''), x+100000, 1700000, cw-200000, 1500000, size=46, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
        rect(slide, x+cw//4, 3400000, cw//2, 40000, P['secondary'])
        txt(slide, s.get('label',''), x+100000, 3550000, cw-200000, 1100000, size=15, color='E0D7FF', align=PP_ALIGN.CENTER)

def impact_v2(slide, data, P):
    """Dark bg + glowing numbers"""
    rect(slide, 0, 0, W, H, P['dark'])
    circle(slide, W//2-600000, H//2-600000, 1200000, P['primary']+'15')
    if data.get('emoji'): txt(slide, data['emoji'], 400000, 300000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1200000, 320000, W-1600000, 700000, size=30, bold=True, color='FFFFFF')
    stats = data.get('stats', [])
    for i, s in enumerate(stats[:3]):
        y = 1600000 + i*1500000
        rect(slide, 400000, y, W-800000, 1200000, P['primary']+'30')
        txt(slide, s.get('number',''), 500000, y+100000, 2500000, 1000000, size=48, bold=True, color=P['secondary'])
        txt(slide, s.get('label',''), 3200000, y+350000, W-3900000, 500000, size=18, color='FFFFFF')

def impact_v3(slide, data, P):
    """Circle stats"""
    rect(slide, 0, 0, W, H, 'FFFFFF')
    rect(slide, 0, 0, W, 1300000, P['primary'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 280000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1100000, 320000, W-1500000, 650000, size=28, bold=True, color='FFFFFF')
    stats = data.get('stats', [])
    cw = W//3
    for i, s in enumerate(stats[:3]):
        cx = i*cw + cw//2
        circle(slide, cx-1000000, 1800000, 2000000, P['light'])
        circle(slide, cx-800000, 2000000, 1600000, P['primary'])
        txt(slide, s.get('number',''), cx-1000000, 2200000, 2000000, 1000000, size=36, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
        txt(slide, s.get('label',''), cx-1200000, 4100000, 2400000, 600000, size=14, color=P['dark'], align=PP_ALIGN.CENTER)

IMPACT_VARIANTS = [impact_v1, impact_v2, impact_v3]

# ═══════════════════════════════════════════════════════════════
# DEMO/FEATURES SLIDE — 3 variants
# ═══════════════════════════════════════════════════════════════
def demo_v1(slide, data, P):
    """2x2 grid"""
    rect(slide, 0, 0, W, 1300000, P['dark'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 260000, 650000, 650000, size=30)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1050000, 300000, W-1400000, 600000, size=28, bold=True, color='FFFFFF')
    features = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    icons = ['🚀','⚡','🔒','📊','🎯','💡']
    cw, ch, gap = (W-600000)//2, (H-1800000)//2, 200000
    for i, f in enumerate(features[:4]):
        x = 300000 + (i%2)*(cw+gap)
        y = 1600000 + (i//2)*(ch+gap)
        rect(slide, x, y, cw, ch, P['light'])
        txt(slide, icons[i%6], x+150000, y+150000, 550000, 550000, size=24)
        txt(slide, f, x+800000, y+180000, cw-950000, ch-360000, size=14, color='374151')

def demo_v2(slide, data, P):
    """Icon list dark"""
    rect(slide, 0, 0, W, H, P['dark'])
    circle(slide, W-1200000, -400000, 2000000, P['primary']+'30')
    if data.get('emoji'): txt(slide, data['emoji'], 400000, 300000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1200000, 320000, W-1600000, 700000, size=30, bold=True, color='FFFFFF')
    features = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    icons = ['🚀','⚡','🔒','📊','🎯','💡','🌟','🔥']
    for i, f in enumerate(features[:5]):
        y = 1400000 + i*1000000
        circle(slide, 400000, y+175000, 650000, P['primary'])
        txt(slide, icons[i%8], 400000, y+175000, 650000, 650000, size=20, align=PP_ALIGN.CENTER)
        txt(slide, f, 1250000, y+150000, W-1700000, 700000, size=17, color='FFFFFF')

def demo_v3(slide, data, P):
    """Colorful cards"""
    rect(slide, 0, 0, W, H, P['bg'])
    rect(slide, 0, 0, W, 1300000, P['primary'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 280000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1100000, 320000, W-1500000, 650000, size=28, bold=True, color='FFFFFF')
    features = data.get('content') or data.get('bullets',[]) or data.get('bullets',[])
    card_colors = [P['primary'], P['secondary'], '374151', '0369a1']
    cw = (W-700000)//2
    ch = (H-1900000)//2
    gap = 180000
    for i, f in enumerate(features[:4]):
        x = 300000 + (i%2)*(cw+gap)
        y = 1500000 + (i//2)*(ch+gap)
        rect(slide, x, y, cw, ch, card_colors[i%4])
        txt(slide, ['💡','⚡','🎯','🔥'][i%4], x+200000, y+200000, 600000, 600000, size=28)
        txt(slide, f, x+200000, y+900000, cw-400000, ch-1100000, size=14, color='FFFFFF')

DEMO_VARIANTS = [demo_v1, demo_v2, demo_v3]

# ═══════════════════════════════════════════════════════════════
# TEAM SLIDE — 3 variants
# ═══════════════════════════════════════════════════════════════
def team_v1(slide, data, P):
    """Avatar circles"""
    rect(slide, 0, 0, W, 1400000, P['primary'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 300000, 700000, 700000, size=32)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1100000, 350000, W-1400000, 600000, size=30, bold=True, color='FFFFFF')
    members = data.get('members', [])
    n = max(1, min(len(members), 4))
    cw = min(2000000, (W-800000)//n)
    gap = 200000
    sx = (W-(n*cw+(n-1)*gap))//2
    for i, m in enumerate(members[:4]):
        x = sx + i*(cw+gap)
        rect(slide, x, 1800000, cw, H-2300000, P['light'])
        av = 900000
        av_x = x+cw//2-av//2
        circle(slide, av_x, 2000000, av, P['primary'])
        initials = ''.join([n[0] for n in m.get('name','TM').split()])[:2].upper()
        txt(slide, initials, av_x, 2000000, av, av, size=26, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
        txt(slide, m.get('name',''), x+100000, 3100000, cw-200000, 500000, size=16, bold=True, color=P['dark'], align=PP_ALIGN.CENTER)
        txt(slide, m.get('role',''), x+100000, 3650000, cw-200000, 400000, size=13, color='374151', align=PP_ALIGN.CENTER)

def team_v2(slide, data, P):
    """Dark cards"""
    rect(slide, 0, 0, W, H, P['dark'])
    if data.get('emoji'): txt(slide, data['emoji'], 400000, 300000, 700000, 700000, size=34)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1200000, 320000, W-1600000, 700000, size=30, bold=True, color='FFFFFF')
    members = data.get('members', [])
    n = max(1, min(len(members), 4))
    cw = min(2000000, (W-800000)//n)
    gap = 200000
    sx = (W-(n*cw+(n-1)*gap))//2
    for i, m in enumerate(members[:4]):
        x = sx + i*(cw+gap)
        rect(slide, x, 1500000, cw, H-2000000, P['primary']+'40')
        rect(slide, x, 1500000, cw, 10000, P['secondary'])
        av = 800000
        circle(slide, x+cw//2-av//2, 1700000, av, P['secondary'])
        initials = ''.join([n[0] for n in m.get('name','TM').split()])[:2].upper()
        txt(slide, initials, x+cw//2-av//2, 1700000, av, av, size=24, bold=True, color=P['dark'], align=PP_ALIGN.CENTER)
        txt(slide, m.get('name',''), x+100000, 2700000, cw-200000, 500000, size=16, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
        txt(slide, m.get('role',''), x+100000, 3250000, cw-200000, 400000, size=13, color=P['secondary'], align=PP_ALIGN.CENTER)

def team_v3(slide, data, P):
    """Horizontal list"""
    rect(slide, 0, 0, W, H, 'FFFFFF')
    rect(slide, 0, 0, 120000, H, P['primary'])
    rect(slide, 0, 0, W, 1300000, P['light'])
    if data.get('emoji'): txt(slide, data['emoji'], 300000, 250000, 650000, 650000, size=30)
    txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 1050000, 280000, W-1400000, 650000, size=28, bold=True, color=P['dark'])
    members = data.get('members', [])
    for i, m in enumerate(members[:4]):
        y = 1500000 + i*1300000
        rect(slide, 300000, y, W-600000, 1100000, P['light'])
        circle(slide, 500000, y+150000, 800000, P['primary'])
        initials = ''.join([n[0] for n in m.get('name','TM').split()])[:2].upper()
        txt(slide, initials, 500000, y+150000, 800000, 800000, size=22, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
        txt(slide, m.get('name',''), 1500000, y+150000, 3000000, 450000, size=18, bold=True, color=P['dark'])
        txt(slide, m.get('role',''), 1500000, y+650000, 3000000, 400000, size=14, color='374151')

TEAM_VARIANTS = [team_v1, team_v2, team_v3]

# ═══════════════════════════════════════════════════════════════
# CLOSING SLIDE — 3 variants
# ═══════════════════════════════════════════════════════════════
def closing_v1(slide, data, P):
    rect(slide, 0, 0, W, H, P['dark'])
    circle(slide, W-2000000, -500000, 2800000, P['primary'])
    circle(slide, -600000, H-1000000, 2000000, P['secondary']+'40')
    if data.get('emoji'): txt(slide, data['emoji'], W//2-400000, 1200000, 800000, 800000, size=48, align=PP_ALIGN.CENTER)
    txt(slide, data.get('title','Thank You'), 457200, 2300000, W-914400, 1000000, size=48, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
    if data.get('subtitle') or data.get('subheading'): txt(slide, data['subtitle'], 457200, 3500000, W-914400, 600000, size=20, color=P['secondary'], italic=True, align=PP_ALIGN.CENTER)
    txt(slide, 'Built with HackMate AI', 457200, H-600000, W-914400, 400000, size=12, color='6b7280', align=PP_ALIGN.CENTER)

def closing_v2(slide, data, P):
    rect(slide, 0, 0, W, H, P['primary'])
    rect(slide, 0, H-2000000, W, 2000000, P['dark'])
    if data.get('emoji'): txt(slide, data['emoji'], W//2-400000, 800000, 800000, 800000, size=56, align=PP_ALIGN.CENTER)
    txt(slide, data.get('title','Thank You'), 457200, 1900000, W-914400, 1200000, size=52, bold=True, color='FFFFFF', align=PP_ALIGN.CENTER)
    if data.get('subtitle') or data.get('subheading'): txt(slide, data['subtitle'], 457200, H-1700000, W-914400, 700000, size=22, color=P['secondary'], italic=True, align=PP_ALIGN.CENTER)

def closing_v3(slide, data, P):
    rect(slide, 0, 0, W, H, 'FFFFFF')
    rect(slide, 0, 0, W, H//2, P['light'])
    circle(slide, W//2-1500000, H//2-1500000, 3000000, P['primary']+'15')
    if data.get('emoji'): txt(slide, data['emoji'], W//2-400000, 800000, 800000, 800000, size=48, align=PP_ALIGN.CENTER)
    txt(slide, data.get('title','Thank You'), 457200, 2000000, W-914400, 1200000, size=52, bold=True, color=P['dark'], align=PP_ALIGN.CENTER)
    if data.get('subtitle') or data.get('subheading'): txt(slide, data['subtitle'], 457200, 3400000, W-914400, 700000, size=20, color=P['primary'], italic=True, align=PP_ALIGN.CENTER)
    rect(slide, W//4, H-800000, W//2, 8000, P['primary'])

CLOSING_VARIANTS = [closing_v1, closing_v2, closing_v3]

# ═══════════════════════════════════════════════════════════════
# MAIN BUILD FUNCTION
# ═══════════════════════════════════════════════════════════════
def build_slide(slide, data, palette, variants_map):
    layout = data.get('layout', 'default')
    variants = variants_map.get(layout)
    if variants:
        # Har baar random variant
        fn = random.choice(variants)
        fn(slide, data, palette)
    else:
        # Default fallback
        rect(slide, 0, 0, 120000, H, palette['primary'])
        rect(slide, 0, 0, W, 1200000, palette['light'])
        txt(slide, data.get('title') or data.get('heading','') or data.get('heading',''), 300000, 280000, W-600000, 650000, size=28, bold=True, color=palette['dark'])
        if data.get('content') or data.get('bullets',[]):
            bullets(slide, data['content'], 300000, 1400000, W-600000, H-1700000, size=16, color='374151')

VARIANTS_MAP = {
    'title': TITLE_VARIANTS,
    'problem': PROBLEM_VARIANTS,
    'solution': SOLUTION_VARIANTS,
    'how': HOW_VARIANTS,
    'tech': TECH_VARIANTS,
    'impact': IMPACT_VARIANTS,
    'demo': DEMO_VARIANTS,
    'team': TEAM_VARIANTS,
    'closing': CLOSING_VARIANTS,
}

# ═══════════════════════════════════════════════════════════════
# AI GENERATION
# ═══════════════════════════════════════════════════════════════
def generate_slides_ai(topic, profile):
    api_key = os.environ.get('ANTHROPIC_API_KEY')
    if not api_key: return fallback_slides(topic, profile)
    try:
        r = requests.post('https://api.anthropic.com/v1/messages',
            headers={'x-api-key': api_key, 'anthropic-version': '2023-06-01', 'Content-Type': 'application/json'},
            json={'model': 'claude-haiku-4-5-20251001', 'max_tokens': 2500, 'messages': [{'role': 'user', 'content': f'''Create presentation for: "{topic}"
Creator: {profile.get("full_name","Student")} | {profile.get("role","Dev")} | {profile.get("college","")}

Return ONLY JSON:
{{"slides":[
  {{"layout":"title","title":"<specific name>","subtitle":"<value prop>","emoji":"🚀"}},
  {{"layout":"problem","title":"The Problem","headline":"<key stat for {topic}>","content":["<pain 1>","<pain 2>","<pain 3>"],"emoji":"😤"}},
  {{"layout":"solution","title":"Our Solution","headline":"<solution>","content":["<feat 1>","<feat 2>","<feat 3>"],"emoji":"💡"}},
  {{"layout":"how","title":"How It Works","steps":["<step 1>","<step 2>","<step 3>","<step 4>"],"emoji":"⚙️"}},
  {{"layout":"tech","title":"Tech Stack","content":["Frontend: <tech>","Backend: <tech>","AI: <tech>","DB: <tech>","Cloud: <tech>"],"emoji":"🛠️"}},
  {{"layout":"impact","title":"Impact","stats":[{{"number":"<stat>","label":"<label>"}},{{"number":"<stat>","label":"<label>"}},{{"number":"<stat>","label":"<label>"}}],"emoji":"📈"}},
  {{"layout":"demo","title":"Key Features","content":["<feat 1>","<feat 2>","<feat 3>","<feat 4>"],"emoji":"✨"}},
  {{"layout":"team","title":"Team","members":[{{"name":"{profile.get("full_name","Lead")}","role":"{profile.get("role","Dev")}"}}],"emoji":"👥"}},
  {{"layout":"closing","title":"Thank You","subtitle":"Let\'s build together","emoji":"🎯"}}
]}}
Rules: SPECIFIC to "{topic}". Bullets max 10 words. Real stats. ONLY JSON.'''}]}, timeout=30)
        text = r.json().get('content',[{}])[0].get('text','')
        clean = text.replace('```json','').replace('```','').strip()
        parsed = json.loads(clean)
        if parsed.get('slides'): return parsed['slides']
    except Exception as e:
        print('AI error:', e)
    return fallback_slides(topic, profile)

def fallback_slides(topic, profile):
    return [
        {'layout':'title','title':topic,'subtitle':'An innovative solution','emoji':'🚀'},
        {'layout':'problem','title':'The Problem','headline':'Critical challenge','content':['Solutions are slow','Real pain points exist','No affordable option'],'emoji':'😤'},
        {'layout':'solution','title':'Our Solution','headline':f'Introducing {topic}','content':['AI-powered','Simple design','Scalable'],'emoji':'💡'},
        {'layout':'how','title':'How It Works','steps':['Input data','AI analyzes','Results ready','Action taken'],'emoji':'⚙️'},
        {'layout':'tech','title':'Tech Stack','content':['Frontend: React','Backend: FastAPI','AI: Claude','DB: PostgreSQL','Cloud: Vercel'],'emoji':'🛠️'},
        {'layout':'impact','title':'Impact','stats':[{'number':'10M+','label':'Users'},{'number':'₹500Cr','label':'Market'},{'number':'80%','label':'Efficiency'}],'emoji':'📈'},
        {'layout':'demo','title':'Key Features','content':['Real-time','Intuitive UI','Offline','Privacy first'],'emoji':'✨'},
        {'layout':'team','title':'Team','members':[{'name':profile.get('full_name','Lead'),'role':profile.get('role','Dev')}],'emoji':'👥'},
        {'layout':'closing','title':'Thank You','subtitle':"Let's build together",'emoji':'🎯'},
    ]

# ═══════════════════════════════════════════════════════════════
# ROUTES
# ═══════════════════════════════════════════════════════════════
@app.route('/generate-ppt', methods=['POST','OPTIONS'])
def generate_ppt():
    if request.method == 'OPTIONS': return '', 200
    try:
        body = request.get_json()
        topic = body.get('topic','My Project')
        theme_name = body.get('theme','purple')
        profile = body.get('profile') or {}
        provided_slides = body.get('slides') or []
        generate_only = body.get('generateOnly', False)

        if generate_only:
            slides = provided_slides or generate_slides_ai(topic, profile)
            return jsonify({'slides': slides})

        slides = provided_slides or generate_slides_ai(topic, profile)

        # Palette — theme se ya random
        palette = THEME_MAP.get(theme_name, random.choice(PALETTES))

        prs = Presentation()
        prs.slide_width = Emu(W)
        prs.slide_height = Emu(H)
        blank = prs.slide_layouts[6]

        for slide_data in slides:
            slide = prs.slides.add_slide(blank)
            build_slide(slide, slide_data, palette, VARIANTS_MAP)

        buf = io.BytesIO()
        prs.save(buf); buf.seek(0)
        filename = (topic[:30]+'.pptx').replace(' ','_')
        return send_file(buf, as_attachment=True, download_name=filename,
                        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status':'ok','service':'HackMate PPT — Random Variants'})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT',5000)))
